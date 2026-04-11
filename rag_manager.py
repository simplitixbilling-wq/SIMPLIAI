"""
RAG (Retrieval-Augmented Generation) Manager
Handles folder processing, chunking, embedding, and vector search
Uses TF-IDF and cosine similarity for efficient retrieval without heavy dependencies
"""

import os
import json
import math
import pickle
import re
import tempfile
from pathlib import Path
from typing import List, Dict, Tuple, Set, Optional
from collections import Counter, defaultdict
import numpy as np
import fitz  # pymupdf — robust PDF extraction (text + OCR)
import pandas as pd

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
except ImportError:
    raise ImportError("sklearn must be installed. Run: pip install scikit-learn")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# Stopwords for query processing (lightweight, no NLTK dependency)
_STOPWORDS = frozenset({
    "a", "an", "the", "is", "are", "was", "were", "be", "been", "being",
    "have", "has", "had", "do", "does", "did", "will", "would", "shall",
    "should", "may", "might", "can", "could", "must", "need",
    "i", "me", "my", "we", "our", "you", "your", "he", "she", "it",
    "they", "them", "their", "this", "that", "these", "those",
    "what", "which", "who", "whom", "when", "where", "how", "why",
    "not", "no", "nor", "but", "or", "and", "if", "then", "so",
    "very", "just", "also", "still", "too", "much", "many", "more",
    "of", "in", "on", "at", "to", "for", "from", "by", "with", "about",
    "into", "through", "during", "before", "after", "above", "below",
    "between", "out", "off", "over", "under", "again", "further",
    "getting", "get", "got", "give", "given", "tell", "told", "please",
    "dont", "im", "its", "ive", "id", "ill", "youre", "youve",
})


class RAGDatabase:
    """Represents a single RAG database with vectors and metadata.
    
    Advanced retrieval inspired by Google NotebookLM:
    - Entity extraction + inverted keyword index (built at index time)
    - Query analysis (intent detection, entity extraction, doc references)
    - Two-pass multi-signal retrieval (TF-IDF -> re-rank)
    - Neighbor context expansion for cross-chunk data
    """
    
    def __init__(self, name: str, embedding_model_name: str = "tfidf"):
        self.name = name
        self.embedding_model_name = embedding_model_name
        self.chunks: List[str] = []
        self.vectorizer = None
        self.embeddings = None
        self.metadata: List[Dict] = []
        self.source_folder: str = ""  # original folder/URL for re-indexing
        # Advanced indexes (rebuilt on load)
        self.keyword_index: Dict[str, List[int]] = {}   # word -> [chunk_ids]
        self.chunk_entities: List[Dict] = []             # per-chunk entities
        
    def add_chunks(self, chunks: List[str], metadata: List[Dict] = None):
        """Add chunks, compute embeddings, and build advanced indexes"""
        if not chunks:
            return
        
        # Create or update vectorizer with all chunks
        all_chunks = self.chunks + chunks
        
        # Use n-gram range (1,2) for better phrase matching
        self.vectorizer = TfidfVectorizer(
            max_features=1500,
            stop_words='english',
            min_df=1,
            ngram_range=(1, 2),
            sublinear_tf=True  # Apply log normalization (BM25-like)
        )
        
        try:
            self.embeddings = self.vectorizer.fit_transform(all_chunks).toarray()
        except:
            # Fallback with simpler settings
            self.vectorizer = TfidfVectorizer(max_features=500, stop_words='english', min_df=1)
            self.embeddings = self.vectorizer.fit_transform(all_chunks).toarray()
        
        # Store chunks and metadata
        self.chunks.extend(chunks)
        if metadata:
            self.metadata.extend(metadata)
        else:
            self.metadata.extend([{"source": f"chunk_{i}"} for i in range(len(chunks))])
        
        # Build advanced indexes
        self._build_keyword_index()
        self._extract_all_entities()
    
    # ------------------------------------------------------------------
    #  INDEXING: Build advanced indexes for multi-signal retrieval
    # ------------------------------------------------------------------

    def _build_keyword_index(self):
        """Build inverted keyword index: word -> [chunk_indices].
        Indexes words, date patterns, and significant numbers for fast lookup."""
        self.keyword_index = defaultdict(list)
        for i, chunk in enumerate(self.chunks):
            # Index meaningful words (3+ chars, not stopwords)
            words = set(re.findall(r'\b[a-zA-Z]{3,}\b', chunk.lower()))
            words -= _STOPWORDS
            for w in words:
                self.keyword_index[w].append(i)
            # Index date patterns (dd/mm/yyyy, dd.mm.yyyy, dd-mm-yyyy)
            for d in re.findall(r'\d{1,2}[./-]\d{1,2}[./-]\d{2,4}', chunk):
                self.keyword_index[d].append(i)
            # Index significant amounts (1,234.00 style)
            for a in re.findall(r'\d[\d,]*\.\d{2}', chunk):
                self.keyword_index[a].append(i)

    def _extract_all_entities(self):
        """Extract entities from every chunk"""
        self.chunk_entities = [self._extract_entities(c) for c in self.chunks]

    @staticmethod
    def _extract_entities(text: str) -> Dict:
        """Extract structured entities: dates, amounts, percentages, key terms"""
        entities = {
            'dates': list(set(re.findall(r'\d{1,2}[./-]\d{1,2}[./-]\d{2,4}', text))),
            'amounts': list(set(re.findall(r'(?:Rs\.?|₹|INR)\s*[\d,]+(?:\.\d+)?', text)
                               + re.findall(r'[\d,]+\.\d{2}', text))),
            'percentages': list(set(re.findall(r'\d+(?:\.\d+)?%', text))),
            'key_terms': [],
        }
        # Capitalized terms = likely proper nouns / org names
        cap_terms = re.findall(r'\b[A-Z][a-zA-Z]{2,}\b', text)
        entities['key_terms'] = list(set(cap_terms))
        return entities

    def _bm25_scores(self, query_terms: List[str], k1: float = 1.5, b: float = 0.75) -> np.ndarray:
        """Pure-Python BM25 scores for all chunks."""
        N = len(self.chunks)
        if N == 0 or not query_terms:
            return np.zeros(N)
        avgdl = sum(len(c.split()) for c in self.chunks) / N
        scores = np.zeros(N)
        for term in query_terms:
            hits = self.keyword_index.get(term)
            if not hits:
                continue
            n_t = len(hits)
            idf = math.log((N - n_t + 0.5) / (n_t + 0.5) + 1.0)
            for idx in hits:
                doc_len = len(self.chunks[idx].split())
                f = self.chunks[idx].lower().count(term)
                tf = (f * (k1 + 1)) / (f + k1 * (1 - b + b * doc_len / max(avgdl, 1)))
                scores[idx] += idf * tf
        mx = scores.max()
        if mx > 0:
            scores /= mx
        return scores

    # ------------------------------------------------------------------
    #  QUERY ANALYSIS: Understand what the user is really asking
    # ------------------------------------------------------------------

    def _analyze_query(self, query: str) -> Dict:
        """Decompose query into intent, key terms, entities, and doc references.
        Inspired by NotebookLM's query understanding layer."""
        query_lower = query.lower()

        # Key terms (remove stopwords)
        words = re.findall(r'\b[a-zA-Z]{2,}\b', query_lower)
        key_terms = [w for w in words if w not in _STOPWORDS]

        # Intent detection
        intent = 'general'
        if any(p in query_lower for p in ['how much', 'how many', 'total', 'amount', 'sum', 'cost', 'price', 'salary', 'ctc', 'tax']):
            intent = 'quantitative'
        elif any(p in query_lower for p in ['when', 'date', 'year', 'month', 'day', 'period', 'time']):
            intent = 'temporal'
        elif any(p in query_lower for p in ['where', 'location', 'address', 'place', 'city', 'office']):
            intent = 'location'
        elif any(p in query_lower for p in ['who', 'name of', 'person', 'employee', 'employer', 'manager']):
            intent = 'entity'
        elif any(p in query_lower for p in ['what is', 'what are', 'what was', 'which']):
            intent = 'factual'
        elif any(p in query_lower for p in ['list', 'all', 'every', 'show me', 'give me', 'summarize', 'summary']):
            intent = 'listing'
        elif any(p in query_lower for p in ['compare', 'difference', 'between', 'vs', 'versus']):
            intent = 'comparison'

        # Extract entities from query itself
        entities = self._extract_entities(query)

        # Identify document/company references by matching query words to filenames
        doc_refs = set()
        explicit_doc_refs = set()
        source_word_map = {}   # source -> words in filename
        for meta in self.metadata:
            src = meta.get('source', '')
            if src and src not in source_word_map:
                name_no_ext = os.path.splitext(src)[0]
                source_word_map[src] = set(
                    w.lower() for w in re.findall(r'[a-zA-Z]{2,}', name_no_ext)
                ) - {'pdf', 'the', 'of', 'is', 'in', 'to', 'and', 'for'}

        for src, name_words in source_word_map.items():
            if set(key_terms) & name_words:
                doc_refs.add(src)

        # Explicit source mention detection: "in <file.pdf>", "from <file>"
        explicit_pattern = re.compile(
            r'(?:in|from|for)\s+([\w\s().,\-]+?\.(?:pdf|docx|csv|txt|xlsx|xls|pptx|ppt|jpg|jpeg|png|webp|bmp|tif|tiff))',
            flags=re.IGNORECASE,
        )
        hinted_sources = [m.strip() for m in explicit_pattern.findall(query)]
        if hinted_sources:
            norm_to_real = {}
            for meta in self.metadata:
                src = meta.get('source', '')
                if src:
                    norm = re.sub(r'\W+', '', src.lower())
                    norm_to_real[norm] = src
            for hint in hinted_sources:
                hnorm = re.sub(r'\W+', '', hint.lower())
                for norm, real in norm_to_real.items():
                    if hnorm and (hnorm in norm or norm in hnorm):
                        explicit_doc_refs.add(real)

        structured_terms = self._extract_structured_query_terms(query)

        requested_fields = []
        field_rules = {
            'pan': [' pan', 'permanent account'],
            'cin': [' cin', 'corporate identity'],
            'passport': ['passport'],
            'ref_no': ['reference number', 'ref no', 'ref no.', 'reference no'],
            'amount': ['amount', 'salary', 'ctc', 'total', 'net pay'],
            'date': ['date', 'dated', 'dob', 'joining'],
        }
        for field, hints in field_rules.items():
            if any(h in query_lower for h in hints):
                requested_fields.append(field)

        return {
            'key_terms': key_terms,
            'intent': intent,
            'entities': entities,
            'doc_refs': list(doc_refs | explicit_doc_refs),
            'explicit_doc_refs': list(explicit_doc_refs),
            'structured_terms': structured_terms,
            'requested_fields': requested_fields,
            'query_lower': query_lower,
        }

    @staticmethod
    def _extract_structured_query_terms(query: str) -> Dict[str, List[str]]:
        """Extract exact structured values from query for precise matching."""
        return {
            'pan': list(set(re.findall(r'\b[A-Z]{5}[0-9]{4}[A-Z]\b', query))),
            'passport': list(set(re.findall(r'\b[A-Z][0-9]{7}\b', query))),
            'cin': list(set(re.findall(r'\b[A-Z0-9]{10,}\b', query))),
            'ids': list(set(re.findall(r'\b[A-Z]{2,}-?[A-Z0-9]{2,}\b', query))),
            'dates': list(set(re.findall(r'\b\d{1,2}[./-]\d{1,2}[./-]\d{2,4}\b', query))),
            'amounts': list(set(re.findall(r'\b\d[\d,]*\.\d{2}\b', query))),
            'roll_ref': list(set(re.findall(r'\b(?:roll|ref|reference)\s*(?:no|number)?\s*[:\-]?\s*([A-Z0-9-]{3,})', query, flags=re.IGNORECASE))),
        }

    @staticmethod
    def _norm(s: str) -> str:
        return re.sub(r'\W+', '', s.lower())

    def _exact_value_score(self, query_info: Dict, chunk_idx: int) -> float:
        """Exact and near-exact value match score for structured terms."""
        terms = query_info.get('structured_terms', {}) or {}
        if not any(terms.values()):
            return 0.0

        chunk = self.chunks[chunk_idx]
        chunk_lower = chunk.lower()
        chunk_norm = self._norm(chunk)

        score = 0.0
        weights = {
            'pan': 0.45,
            'passport': 0.45,
            'cin': 0.35,
            'ids': 0.30,
            'roll_ref': 0.30,
            'dates': 0.25,
            'amounts': 0.25,
        }

        for key, vals in terms.items():
            w = weights.get(key, 0.2)
            for val in vals:
                v = val.strip()
                if not v:
                    continue
                if v.lower() in chunk_lower:
                    score += w
                    continue
                # OCR/noisy spacing tolerant match
                if self._norm(v) and self._norm(v) in chunk_norm:
                    score += w * 0.9

        return min(score, 1.2)

    def _field_pattern_score(self, query_info: Dict, chunk_idx: int) -> float:
        """Boost chunks that likely contain the requested field even if exact value isn't in query."""
        fields = query_info.get('requested_fields', []) or []
        if not fields:
            return 0.0

        text = self.chunks[chunk_idx]
        lower = text.lower()
        score = 0.0

        for f in fields:
            if f == 'pan':
                if 'pan' in lower:
                    score += 0.5
                if re.search(r'\b[A-Z]{5}[0-9]{4}[A-Z]\b', text):
                    score += 0.6
            elif f == 'cin':
                if 'cin' in lower:
                    score += 0.5
                if re.search(r'\b[A-Z0-9]{12,}\b', text):
                    score += 0.4
            elif f == 'passport':
                if 'passport' in lower:
                    score += 0.6
                if re.search(r'\b[A-Z][0-9]{7}\b', text):
                    score += 0.5
            elif f == 'ref_no':
                if 'ref no' in lower or 'reference' in lower:
                    score += 0.7
                if re.search(r'\b[A-Z0-9-]{5,}\b', text):
                    score += 0.2
            elif f == 'amount':
                if any(k in lower for k in ['amount', 'salary', 'ctc', 'net pay', 'total']):
                    score += 0.5
                if re.search(r'\b\d[\d,]*\.\d{2}\b', text):
                    score += 0.4
            elif f == 'date':
                if 'date' in lower or 'dated' in lower:
                    score += 0.4
                if re.search(r'\b\d{1,2}[./-]\d{1,2}[./-]\d{2,4}\b', text):
                    score += 0.4

        return min(score, 1.2)

    # ------------------------------------------------------------------
    #  SCORING SIGNALS for re-ranking
    # ------------------------------------------------------------------

    def _keyword_overlap_score(self, query_info: Dict, chunk_idx: int) -> float:
        """Exact keyword overlap score. Includes multi-word phrase matching."""
        key_terms = query_info['key_terms']
        if not key_terms:
            return 0.0

        chunk_lower = self.chunks[chunk_idx].lower()
        # Single-word matches
        matches = sum(1 for t in key_terms if t in chunk_lower)
        base = matches / len(key_terms)

        # Multi-word phrase bonus (2- and 3-grams from query)
        phrase_bonus = 0.0
        words = query_info['query_lower'].split()
        for n in (3, 2):
            for i in range(len(words) - n + 1):
                phrase = ' '.join(words[i:i+n])
                if phrase in chunk_lower and not all(w in _STOPWORDS for w in phrase.split()):
                    phrase_bonus += 0.12 * n
        return base * 0.5 + min(phrase_bonus, 0.5)

    def _entity_match_score(self, query_info: Dict, chunk_idx: int) -> float:
        """Score based on entity overlap between query and chunk."""
        if chunk_idx >= len(self.chunk_entities):
            return 0.0
        q_ents = query_info['entities']
        c_ents = self.chunk_entities[chunk_idx]
        score = 0.0
        # Date match
        for qd in q_ents.get('dates', []):
            if qd in c_ents.get('dates', []):
                score += 0.3
        # Amount match
        for qa in q_ents.get('amounts', []):
            if qa in c_ents.get('amounts', []):
                score += 0.2
        # Key-term match (proper nouns)
        q_terms = set(t.lower() for t in q_ents.get('key_terms', []))
        c_terms = set(t.lower() for t in c_ents.get('key_terms', []))
        overlap = q_terms & c_terms
        if overlap:
            score += 0.15 * min(len(overlap), 3)
        return score

    def _intent_content_score(self, intent: str, chunk_idx: int) -> float:
        """Score a chunk based on how well its content matches the query intent."""
        chunk = self.chunks[chunk_idx]
        if intent == 'quantitative':
            num_ratio = sum(1 for c in chunk if c.isdigit()) / max(len(chunk), 1)
            return min(num_ratio * 3, 0.15)
        elif intent == 'temporal':
            date_cnt = len(re.findall(r'\d{1,2}[./-]\d{1,2}[./-]\d{2,4}', chunk))
            return min(date_cnt * 0.04, 0.15)
        elif intent in ('entity', 'factual'):
            words = chunk.split()
            if words:
                cap = sum(1 for w in words if w and w[0].isupper()) / len(words)
                return min(cap * 0.2, 0.10)
        elif intent == 'listing':
            # Prefer chunks with list-like structure
            lines = chunk.split('\n')
            if len(lines) > 3:
                return 0.10
        return 0.0

    # ------------------------------------------------------------------
    #  RETRIEVE: Two-pass multi-signal retrieval
    # ------------------------------------------------------------------

    def retrieve(self, query: str, k: int = 5, source_filter: Optional[List[str]] = None) -> List[Tuple[str, float]]:
        """
        Two-pass multi-signal retrieval (inspired by NotebookLM).

        Pass 1 – Broad candidate collection:
            • TF-IDF cosine similarity (semantic-ish via bigrams)
            • Inverted keyword index hits (exact lexical match)
            → Union of top candidates from both signals

        Pass 2 – Multi-signal re-ranking:
            Signal 1: TF-IDF similarity            (weight 0.30)
            Signal 2: Exact keyword/phrase overlap  (weight 0.25)
            Signal 3: Entity overlap                (weight 0.15)
            Signal 4: Document-name match           (weight 0.20)
            Signal 5: Intent-content match          (weight 0.10)

        Post-ranking:
            • Neighbor context boost (±1 chunks)
            • Deduplication

        Returns list of (chunk, score) tuples.
        """
        if self.vectorizer is None or not self.chunks:
            return []

        try:
            # --- Query Analysis ---
            query_info = self._analyze_query(query)

            # --- PASS 1: BROAD CANDIDATE COLLECTION ---
            # 1a. TF-IDF similarity
            query_vec = self.vectorizer.transform([query]).toarray()
            tfidf_sims = cosine_similarity(query_vec, self.embeddings)[0]

            # 1b. Keyword index lookup (fast exact matching)
            keyword_hits = np.zeros(len(self.chunks))
            for term in query_info['key_terms']:
                if term in self.keyword_index:
                    for cid in self.keyword_index[term]:
                        keyword_hits[cid] += 1.0 / max(len(query_info['key_terms']), 1)

            # 1c. BM25 scores
            bm25 = self._bm25_scores(query_info['key_terms'])

            # Merge signals for candidate selection
            candidate_scores = tfidf_sims + keyword_hits * 0.3 + bm25 * 0.2
            n_candidates = min(25, len(self.chunks))
            candidate_set = set(int(i) for i in np.argsort(candidate_scores)[::-1][:n_candidates])

            allowed_sources = set(source_filter or [])
            if query_info.get('explicit_doc_refs'):
                allowed_sources |= set(query_info['explicit_doc_refs'])

            if allowed_sources:
                scoped = set()
                for i, meta in enumerate(self.metadata):
                    if meta.get('source', '') in allowed_sources:
                        scoped.add(i)
                if scoped:
                    # Hard source scope when explicitly constrained.
                    candidate_set = (candidate_set & scoped) or scoped

            # Also force-include chunks from referenced documents
            if query_info['doc_refs']:
                for i, meta in enumerate(self.metadata):
                    if meta.get('source', '') in query_info['doc_refs']:
                        candidate_set.add(i)

            # --- PASS 2: MULTI-SIGNAL RE-RANKING ---
            final_scores = {}
            for idx in candidate_set:
                if idx >= len(self.chunks):
                    continue

                s = 0.0
                # Signal 1: TF-IDF (0.25)
                s += tfidf_sims[idx] * 0.25
                # Signal 2: Keyword/phrase overlap (0.20)
                s += self._keyword_overlap_score(query_info, idx) * 0.20
                # Signal 3: Entity overlap (0.15)
                s += self._entity_match_score(query_info, idx) * 0.15
                # Signal 4: Document-name match (0.20)
                if query_info['doc_refs']:
                    src = self.metadata[idx].get('source', '') if idx < len(self.metadata) else ''
                    if src in query_info['doc_refs']:
                        if query_info.get('explicit_doc_refs'):
                            s += 0.35
                        else:
                            s += 0.20
                    else:
                        if query_info.get('explicit_doc_refs'):
                            s *= 0.35
                        else:
                            s *= 0.7   # mild penalty
                # Signal 5: Intent-content match (0.10)
                s += self._intent_content_score(query_info['intent'], idx) * 0.10
                # Signal 6: BM25 (0.10)
                s += float(bm25[idx]) * 0.10
                # Signal 7: exact structured-term value match (0.20..)
                s += self._exact_value_score(query_info, idx) * 0.20
                # Signal 8: field-intent pattern score
                s += self._field_pattern_score(query_info, idx) * 0.12

                final_scores[idx] = s

            # --- NEIGHBOR CONTEXT BOOST ---
            neighbor_add = {}
            for idx, score in list(final_scores.items()):
                if score > 0.08:
                    boost = score * 0.25
                    for nb in (idx - 1, idx + 1):
                        if 0 <= nb < len(self.chunks):
                            if nb not in final_scores:
                                final_scores[nb] = tfidf_sims[nb] * 0.15
                            neighbor_add[nb] = max(neighbor_add.get(nb, 0), boost)
            for nb, boost in neighbor_add.items():
                final_scores[nb] = final_scores.get(nb, 0) + boost

            # --- SELECT TOP-K WITH DEDUP ---
            sorted_ids = sorted(final_scores, key=final_scores.get, reverse=True)

            results = []
            seen = set()
            for idx in sorted_ids:
                if len(results) >= k:
                    break
                if final_scores[idx] < 0.01:
                    continue
                chunk = self.chunks[idx]
                sig = chunk[:120]
                if sig not in seen:
                    seen.add(sig)
                    results.append((chunk, float(final_scores[idx])))

            return results
        except Exception as _exc:
            print(f"[RAG] Retrieve error in '{self.name}': {_exc}")
            return []
    
    def get_chunk_source(self, chunk_idx: int) -> str:
        """Get the source file for a given chunk index by tracing --- File: markers"""
        if chunk_idx < 0 or chunk_idx >= len(self.chunks):
            return ""
        
        # Check metadata first (new format)
        if chunk_idx < len(self.metadata):
            src = self.metadata[chunk_idx].get("source", "")
            if src and not src.startswith(("chunk_", self.name + "_chunk")):
                return src
        
        # Check if chunk itself has a file marker - use LAST match (most recent file)
        matches = re.findall(r'--- File: (.+?) ---', self.chunks[chunk_idx])
        if matches:
            return matches[-1]  # Last file marker is the current file
        
        # Trace backwards to find nearest --- File: marker
        for i in range(chunk_idx - 1, -1, -1):
            matches = re.findall(r'--- File: (.+?) ---', self.chunks[i])
            if matches:
                return matches[-1]  # Last marker in that chunk
        
        return ""

    def save(self, directory: str):
        """Save RAG database to disk"""
        os.makedirs(directory, exist_ok=True)
        
        # Save metadata
        metadata_path = os.path.join(directory, "metadata.json")
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump({
                "name": self.name,
                "embedding_model": self.embedding_model_name,
                "num_chunks": len(self.chunks),
                "chunk_metadata": self.metadata,
                "source_folder": self.source_folder
            }, f, indent=2, ensure_ascii=True)
        
        # Save chunks
        chunks_path = os.path.join(directory, "chunks.json")
        with open(chunks_path, "w", encoding="utf-8") as f:
            json.dump(self.chunks, f, ensure_ascii=True)
        
        # Save vectorizer
        if self.vectorizer is not None:
            vectorizer_path = os.path.join(directory, "vectorizer.pkl")
            with open(vectorizer_path, "wb") as f:
                pickle.dump(self.vectorizer, f)
            
            embeddings_path = os.path.join(directory, "embeddings.pkl")
            with open(embeddings_path, "wb") as f:
                pickle.dump(self.embeddings, f)

        # Cache advanced indexes so they don't need to be rebuilt on every load
        indexes_path = os.path.join(directory, "indexes.pkl")
        with open(indexes_path, "wb") as f:
            pickle.dump({
                "keyword_index": self.keyword_index,
                "chunk_entities": self.chunk_entities
            }, f)
    
    @staticmethod
    def load(directory: str) -> 'RAGDatabase':
        """Load RAG database from disk"""
        metadata_path = os.path.join(directory, "metadata.json")
        
        with open(metadata_path, "r") as f:
            metadata = json.load(f)
        
        db = RAGDatabase(
            name=metadata["name"],
            embedding_model_name=metadata["embedding_model"]
        )
        
        # Load chunks
        chunks_path = os.path.join(directory, "chunks.json")
        if os.path.exists(chunks_path):
            with open(chunks_path, "r", encoding="utf-8") as f:
                db.chunks = json.load(f)
        
        # Load vectorizer and embeddings
        vectorizer_path = os.path.join(directory, "vectorizer.pkl")
        if os.path.exists(vectorizer_path):
            with open(vectorizer_path, "rb") as f:
                db.vectorizer = pickle.load(f)
            
            embeddings_path = os.path.join(directory, "embeddings.pkl")
            if os.path.exists(embeddings_path):
                with open(embeddings_path, "rb") as f:
                    db.embeddings = pickle.load(f)
        
        db.metadata = metadata.get("chunk_metadata", [])
        db.source_folder = metadata.get("source_folder", "")
        # Fallback: infer source_folder from first chunk's metadata
        if not db.source_folder and db.metadata:
            db.source_folder = db.metadata[0].get("folder", "")
        
        # Load advanced indexes from cache if available; otherwise rebuild
        indexes_path = os.path.join(directory, "indexes.pkl")
        if os.path.exists(indexes_path):
            try:
                with open(indexes_path, "rb") as f:
                    cached = pickle.load(f)
                db.keyword_index = cached.get("keyword_index", {})
                db.chunk_entities = cached.get("chunk_entities", [])
            except Exception:
                # Corrupted cache — rebuild
                if db.chunks:
                    db._build_keyword_index()
                    db._extract_all_entities()
        elif db.chunks:
            db._build_keyword_index()
            db._extract_all_entities()
        
        return db


class RAGManager:
    """Manages multiple RAG databases"""
    
    def __init__(self, base_directory: str = "rag_databases"):
        self.base_directory = base_directory
        self.databases: Dict[str, RAGDatabase] = {}
        self._rapidocr_engine = None
        self._md_snapshot_cache: Dict[str, Dict] = {}
        os.makedirs(base_directory, exist_ok=True)
        self._load_all_databases()

    def _get_rapidocr_engine(self):
        if self._rapidocr_engine is not None:
            return self._rapidocr_engine
        try:
            from rapidocr_onnxruntime import RapidOCR
            self._rapidocr_engine = RapidOCR()
            return self._rapidocr_engine
        except Exception:
            return None

    def _ocr_image_with_rapidocr(self, image_path: str) -> str:
        engine = self._get_rapidocr_engine()
        if engine is None:
            return ""

        try:
            from PIL import Image, ImageOps, ImageEnhance, ImageFilter

            def _collect_lines(ocr_result):
                lines = []
                for item in ocr_result or []:
                    if len(item) >= 2 and isinstance(item[1], str):
                        txt = item[1].strip()
                        if not txt:
                            continue
                        conf = 0.0
                        if len(item) >= 3:
                            try:
                                conf = float(item[2])
                            except Exception:
                                conf = 0.0
                        lines.append((conf, txt))
                return lines

            best_text = ""
            best_score = -1.0

            # Pass 1: original image
            result, _elapsed = engine(image_path)
            lines = _collect_lines(result)
            if lines:
                score = sum(conf for conf, _txt in lines)
                text = "\n".join(txt for _conf, txt in lines)
                if text.strip() and score > best_score:
                    best_text, best_score = text, score

            # Pass 2: enhanced variants
            with tempfile.TemporaryDirectory() as temp_dir:
                src = Image.open(image_path).convert("L")
                variants = []

                up = src.resize((src.width * 2, src.height * 2))
                up = ImageOps.autocontrast(up)
                variants.append(up)

                sharp = ImageEnhance.Contrast(src).enhance(2.2)
                sharp = sharp.filter(ImageFilter.SHARPEN)
                variants.append(sharp)

                binary = src.point(lambda p: 255 if p > 160 else 0, mode="1").convert("L")
                variants.append(binary)

                for idx, image in enumerate(variants):
                    temp_path = os.path.join(temp_dir, f"ocr_variant_{idx}.png")
                    image.save(temp_path)
                    cand_result, _cand_elapsed = engine(temp_path)
                    cand_lines = _collect_lines(cand_result)
                    if not cand_lines:
                        continue
                    cand_score = sum(conf for conf, _txt in cand_lines)
                    cand_text = "\n".join(txt for _conf, txt in cand_lines)
                    if cand_text.strip() and cand_score > best_score:
                        best_text, best_score = cand_text, cand_score

            return best_text.strip()
        except Exception:
            return ""

    def _ocr_pdf_page_with_rapidocr(self, page) -> str:
        try:
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            with tempfile.TemporaryDirectory() as temp_dir:
                image_path = os.path.join(temp_dir, "page.png")
                pix.save(image_path)
                return self._ocr_image_with_rapidocr(image_path)
        except Exception:
            return ""
    
    def _load_all_databases(self):
        """Load all saved RAG databases from disk"""
        if not os.path.exists(self.base_directory):
            return
        
        for db_name in os.listdir(self.base_directory):
            db_path = os.path.join(self.base_directory, db_name)
            if os.path.isdir(db_path):
                try:
                    self.databases[db_name] = RAGDatabase.load(db_path)
                    # Backfill markdown snapshot for older databases created before this feature.
                    md_path = self.get_markdown_path(db_name)
                    if not os.path.exists(md_path):
                        self.export_database_markdown(db_name)
                except Exception as e:
                    print(f"Error loading RAG database {db_name}: {e}")

    def get_markdown_path(self, rag_name: str) -> str:
        """Return canonical markdown snapshot path for a RAG database."""
        return os.path.join(self.base_directory, rag_name, "knowledge.md")

    def export_database_markdown(self, rag_name: str, max_chunk_chars: int = 1200) -> str:
        """Write a markdown snapshot for quick human inspection and lightweight text lookup.

        Returns the markdown file path.
        """
        if rag_name not in self.databases:
            raise ValueError(f"Database '{rag_name}' not found")

        db = self.databases[rag_name]
        md_path = self.get_markdown_path(rag_name)
        os.makedirs(os.path.dirname(md_path), exist_ok=True)

        lines = []
        lines.append(f"# RAG Knowledge Snapshot: {rag_name}")
        lines.append("")
        lines.append(f"- Total chunks: {len(db.chunks)}")
        lines.append(f"- Embedding model: {db.embedding_model_name}")
        if db.source_folder:
            lines.append(f"- Source: {db.source_folder}")
        lines.append("")

        # Build compact chunk sections grouped by source file.
        by_source = defaultdict(list)
        for idx, chunk in enumerate(db.chunks):
            source = ""
            if idx < len(db.metadata):
                source = db.metadata[idx].get("source", "")
            if not source:
                source = db.get_chunk_source(idx) or "unknown"
            by_source[source].append((idx, chunk))

        for source in sorted(by_source.keys()):
            lines.append(f"## Source: {source}")
            lines.append("")
            for idx, chunk in by_source[source]:
                lines.append(f"### Chunk {idx}")
                lines.append("")
                cleaned = re.sub(r'^---\s*File:\s*.+?\s*---\s*', '', chunk, flags=re.IGNORECASE)
                cleaned = cleaned.strip()
                if len(cleaned) > max_chunk_chars:
                    cleaned = cleaned[:max_chunk_chars].rstrip() + " ..."
                lines.append(cleaned if cleaned else "(empty chunk)")
                lines.append("")

        with open(md_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines).strip() + "\n")

        return md_path
    
    def create_from_folder(self, folder_path: str, rag_name: str, chunk_size: int = 512, chunk_overlap: int = 100, progress_callback=None) -> RAGDatabase:
        """
        Create RAG database from a folder of documents
        
        Args:
            folder_path: Path to folder containing documents
            rag_name: Name for the RAG database
            chunk_size: Size of each text chunk
            chunk_overlap: Overlap between chunks
            progress_callback: Optional callable(file_name, file_index, total_files)
        
        Returns:
            RAGDatabase instance
        """
        if not os.path.isdir(folder_path):
            raise ValueError(f"Folder not found: {folder_path}")
        
        # Check if name already exists
        if rag_name in self.databases:
            raise ValueError(f"RAG database '{rag_name}' already exists")
        
        # Extract text from all files in folder (per-file)
        file_texts = self._extract_texts_from_folder(folder_path, progress_callback=progress_callback)
        
        if not file_texts:
            raise ValueError("No text content found in folder")
        
        # Create chunks with per-file source metadata
        all_chunks = []
        all_metadata = []
        
        for file_name, text in file_texts:
            if not text.strip():
                continue
            chunks = self._chunk_text(text, chunk_size, chunk_overlap)
            # Keep short OCR outputs (e.g., scanned IDs) indexable.
            if not chunks and text.strip():
                chunks = [text.strip()]
            for i, chunk in enumerate(chunks):
                # Prepend file marker so retrieval can cite the source doc
                tagged_chunk = f"--- File: {file_name} ---\n{chunk}"
                all_chunks.append(tagged_chunk)
                all_metadata.append({
                    "source": file_name,
                    "chunk_index": i,
                    "folder": folder_path
                })
        
        if not all_chunks:
            raise ValueError("No text content found in folder")
        
        # Create RAG database
        db = RAGDatabase(name=rag_name)
        db.source_folder = folder_path  # store for re-indexing
        db.add_chunks(all_chunks, all_metadata)
        
        # Save to disk
        db_path = os.path.join(self.base_directory, rag_name)
        db.save(db_path)
        
        # Store in memory
        self.databases[rag_name] = db

        # Export markdown snapshot for quick browsing / lightweight grep workflows.
        self.export_database_markdown(rag_name)
        
        return db
    
    def reindex_database(self, rag_name: str, chunk_size: int = 512, chunk_overlap: int = 100, progress_callback=None) -> 'RAGDatabase':
        """Re-process documents from the original source folder."""
        if rag_name not in self.databases:
            raise ValueError(f"Database '{rag_name}' not found")
        
        db = self.databases[rag_name]
        source_folder = db.source_folder or (db.metadata[0].get("folder", "") if db.metadata else "")
        
        if not source_folder or not os.path.isdir(source_folder):
            raise ValueError(f"Source folder not found for '{rag_name}'. Cannot re-index.")
        
        # Delete existing entry
        del self.databases[rag_name]
        db_path = os.path.join(self.base_directory, rag_name)
        if os.path.exists(db_path):
            import shutil
            shutil.rmtree(db_path)
        
        return self.create_from_folder(source_folder, rag_name, chunk_size, chunk_overlap, progress_callback=progress_callback)
    
    def create_from_url(self, url: str, rag_name: str, chunk_size: int = 512, chunk_overlap: int = 100) -> 'RAGDatabase':
        """Create RAG database from a web page URL."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 is required: pip install beautifulsoup4")
        
        if rag_name in self.databases:
            raise ValueError(f"RAG database '{rag_name}' already exists")
        
        # Fetch page with retry
        import requests as _requests
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
        last_error = None
        html = None
        for _attempt in range(3):
            try:
                resp = _requests.get(url, headers=headers, timeout=20)
                resp.raise_for_status()
                html = resp.content.decode('utf-8', errors='ignore')
                break
            except Exception as exc:
                last_error = exc
        if html is None:
            raise ValueError(f"Failed to fetch URL after 3 attempts: {last_error}")
        
        soup = BeautifulSoup(html, 'html.parser')
        for tag in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
            tag.decompose()
        text = soup.get_text(separator='\n', strip=True)
        
        if not text.strip():
            raise ValueError("No text content found at URL")
        
        from urllib.parse import urlparse
        parsed = urlparse(url)
        page_name = (parsed.netloc + parsed.path).rstrip('/') or "webpage"
        
        chunks = self._chunk_text(text, chunk_size, chunk_overlap)
        tagged_chunks = [f"--- File: {page_name} ---\n{c}" for c in chunks]
        metadata = [{"source": page_name, "chunk_index": i, "folder": url} for i in range(len(chunks))]
        
        db = RAGDatabase(name=rag_name)
        db.source_folder = url
        db.add_chunks(tagged_chunks, metadata)
        
        db_path = os.path.join(self.base_directory, rag_name)
        db.save(db_path)
        self.databases[rag_name] = db

        # Export markdown snapshot for quick browsing / lightweight grep workflows.
        self.export_database_markdown(rag_name)
        
        return db
    
    def _extract_texts_from_folder(self, folder_path: str, progress_callback=None) -> List[Tuple[str, str]]:
        """Extract text from all supported files in a folder.
        Returns list of (filename, text) tuples."""
        file_texts = []
        supported_extensions = {
            ".txt", ".pdf", ".csv", ".xlsx", ".xls", ".docx", ".pptx", ".ppt",
            ".jpg", ".jpeg", ".png", ".bmp", ".webp", ".tif", ".tiff"
        }
        
        all_files = [f for f in sorted(Path(folder_path).rglob("*"))
                     if f.suffix.lower() in supported_extensions]
        total = len(all_files)
        
        for idx, file_path in enumerate(all_files):
            if progress_callback:
                try:
                    progress_callback(file_path.name, idx, total)
                except Exception:
                    pass
            try:
                text = self._extract_text_from_file(str(file_path))
                if text and text.strip():
                    file_texts.append((file_path.name, text))
            except Exception as e:
                print(f"Error processing file {file_path}: {e}")
                continue
        
        if progress_callback and total > 0:
            try:
                progress_callback("Done", total, total)
            except Exception:
                pass
        
        return file_texts
    
    def _extract_text_from_folder(self, folder_path: str) -> str:
        """Legacy: Extract text from all supported files as single string"""
        file_texts = self._extract_texts_from_folder(folder_path)
        all_text = ""
        for fname, text in file_texts:
            all_text += f"\n\n--- File: {fname} ---\n\n{text}"
        return all_text
    
    def _extract_text_from_file(self, file_path: str) -> str:
        """Extract text from a single file"""
        ext = Path(file_path).suffix.lower()
        
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        
        elif ext == ".pdf":
            text = ""
            ocr_used = False
            doc = fitz.open(file_path)
            for page in doc:
                page_text = page.get_text("text")
                if page_text and page_text.strip():
                    text += page_text + "\n"
                else:
                    # Scanned page — OCR fallback (PyMuPDF/Tesseract first)
                    ocr_used = True
                    try:
                        ocr_text = page.get_textpage_ocr().extractText()
                        if ocr_text and ocr_text.strip():
                            text += ocr_text + "\n"
                        else:
                            raise RuntimeError("empty tesseract output")
                    except Exception:
                        # Tesseract unavailable — fallback to rapidocr_onnxruntime
                        ocr_text = self._ocr_pdf_page_with_rapidocr(page)
                        if ocr_text:
                            text += ocr_text + "\n"
            doc.close()
            if ocr_used:
                print(f"[RAG] Scanned PDF detected — OCR applied to: {os.path.basename(file_path)}")
            return text

        elif ext in (".jpg", ".jpeg", ".png", ".bmp", ".webp", ".tif", ".tiff"):
            return self._ocr_image_with_rapidocr(file_path)
        
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            return df.to_string()
        
        elif ext in (".xlsx", ".xls"):
            # Read all sheets
            xls = pd.ExcelFile(file_path)
            texts = []
            for sheet in xls.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                texts.append(f"[Sheet: {sheet}]\n{df.to_string()}")
            return "\n\n".join(texts)
        
        elif ext == ".docx":
            from docx import Document
            text = ""
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text += row_text + "\n"
            return text
        
        elif ext in (".pptx", ".ppt"):
            if not PPTX_AVAILABLE:
                print(f"Skipping {file_path}: python-pptx not installed")
                return ""
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"[Slide {slide_num}]\n"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_text += para.text.strip() + "\n"
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                slide_text += row_text + "\n"
                if slide_text.strip() != f"[Slide {slide_num}]":
                    text += slide_text + "\n"
            return text
        
        return ""
    
    def _chunk_text(self, text: str, chunk_size: int = 512, chunk_overlap: int = 100) -> List[str]:
        """
        Split text into overlapping chunks using sentence-aware boundaries.
        """
        # Split on sentence boundaries (period/newline) rather than only newlines
        sentences = re.split(r'(?<=[.!?\n])\s+', text)
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            # If adding this sentence exceeds chunk size, finalize current chunk
            if len(current_chunk) + len(sentence) > chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                # Keep overlap: last ~15 words for better context continuity
                words = current_chunk.split()
                overlap_words = min(15, len(words) // 3)
                overlap_text = " ".join(words[-overlap_words:]) if overlap_words > 0 else ""
                current_chunk = overlap_text + " " + sentence if overlap_text else sentence
            else:
                current_chunk += (" " + sentence) if current_chunk else sentence
        
        # Add final chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        # Filter out very small chunks, but always keep lines with financial entities
        def _should_keep(c: str) -> bool:
            if len(c.strip()) > 30:
                return True
            return bool(re.search(
                r'\d[\d,]*\.\d{2}|[\u20b9$€£]|total|gstin|pan|invoice|date|ref',
                c, re.IGNORECASE
            ))
        return [c for c in chunks if _should_keep(c)]

    def _load_markdown_snapshot_entries(self, rag_name: str) -> List[Dict[str, str]]:
        """Parse knowledge.md into searchable entries and cache by mtime."""
        md_path = self.get_markdown_path(rag_name)
        if not os.path.exists(md_path):
            return []

        mtime = os.path.getmtime(md_path)
        cached = self._md_snapshot_cache.get(rag_name)
        if cached and cached.get("mtime") == mtime:
            return cached.get("entries", [])

        entries: List[Dict[str, str]] = []
        current_source = "unknown"
        current_chunk = ""
        current_lines: List[str] = []

        def flush_entry():
            nonlocal current_chunk, current_lines
            text = "\n".join(current_lines).strip()
            if current_chunk and text:
                entries.append({
                    "source": current_source,
                    "chunk": current_chunk,
                    "text": text,
                    "lower": text.lower(),
                })
            current_chunk = ""
            current_lines = []

        with open(md_path, "r", encoding="utf-8", errors="ignore") as f:
            for raw in f:
                line = raw.rstrip("\n")
                if line.startswith("## Source: "):
                    flush_entry()
                    current_source = line[len("## Source: "):].strip() or "unknown"
                    continue
                if line.startswith("### Chunk "):
                    flush_entry()
                    current_chunk = line[len("### Chunk "):].strip()
                    continue
                if current_chunk:
                    current_lines.append(line)

        flush_entry()

        self._md_snapshot_cache[rag_name] = {
            "mtime": mtime,
            "entries": entries,
        }
        return entries

    @staticmethod
    def _query_terms(query: str) -> List[str]:
        words = re.findall(r"\b[a-zA-Z0-9]{2,}\b", query.lower())
        return [w for w in words if w not in _STOPWORDS]

    def _search_markdown_snapshot(self, rag_name: str, query: str, top_n: int = 8) -> List[Tuple[str, float]]:
        """Fast lexical pre-search over knowledge.md snapshot."""
        entries = self._load_markdown_snapshot_entries(rag_name)
        if not entries:
            return []

        q_lower = query.lower()
        terms = self._query_terms(query)
        scored: List[Tuple[float, Dict[str, str]]] = []

        for e in entries:
            text_l = e["lower"]
            term_hits = sum(1 for t in terms if t in text_l)
            if terms:
                kw_score = term_hits / len(terms)
            else:
                kw_score = 0.0

            phrase_score = 0.5 if q_lower and q_lower in text_l else 0.0
            # Keep entries with at least weak lexical overlap.
            score = kw_score * 0.8 + phrase_score
            if score > 0.08:
                scored.append((score, e))

        scored.sort(key=lambda x: x[0], reverse=True)
        out: List[Tuple[str, float]] = []
        for score, e in scored[:top_n]:
            text = e["text"]
            # Bound markdown snippet length before forwarding to prompt context.
            if len(text) > 1200:
                text = text[:1200].rstrip() + " ..."
            chunk = f"--- File: {e['source']} ---\n{text}"
            out.append((chunk, float(score)))
        return out
    
    def retrieve(self, rag_name: str, query: str, k: int = 5, source_filter: Optional[List[str]] = None) -> List[Tuple[str, float]]:
        """
        Hybrid retrieve:
        1) Fast markdown lexical pre-search
        2) Vector/semantic retrieval
        3) Merge and re-rank
        """
        if rag_name not in self.databases:
            return []

        # Pull deeper candidate pools for better merge quality.
        vec_candidates = self.databases[rag_name].retrieve(
            query,
            max(k * 3, 8),
            source_filter=source_filter,
        )
        md_candidates = self._search_markdown_snapshot(
            rag_name,
            query,
            top_n=max(k * 2, 6),
        )

        if not md_candidates:
            return vec_candidates[:k]
        if not vec_candidates:
            return md_candidates[:k]

        # Merge by normalized content key; preserve strongest signal.
        merged: Dict[str, Tuple[str, float]] = {}

        def _key(text: str) -> str:
            return re.sub(r"\s+", " ", text.strip().lower())[:500]

        for chunk, score in vec_candidates:
            key = _key(chunk)
            merged[key] = (chunk, float(score))

        for chunk, score in md_candidates:
            key = _key(chunk)
            # Markdown is lexical and fast; use modest blend weight.
            md_weighted = 0.25 + float(score) * 0.35
            if key in merged:
                old_chunk, old_score = merged[key]
                merged[key] = (old_chunk, max(old_score, md_weighted))
            else:
                merged[key] = (chunk, md_weighted)

        ranked = sorted(merged.values(), key=lambda x: x[1], reverse=True)
        return ranked[:k]
    
    def get_database_info(self, rag_name: str) -> Dict:
        """Get information about a RAG database"""
        if rag_name not in self.databases:
            return None
        
        db = self.databases[rag_name]
        return {
            "name": rag_name,
            "num_chunks": len(db.chunks),
            "embedding_model": db.embedding_model_name,
            "status": "Ready"
        }
    
    def list_databases(self) -> List[str]:
        """List all available RAG databases"""
        return list(self.databases.keys())
    
    def delete_database(self, rag_name: str) -> bool:
        """Delete a RAG database"""
        if rag_name not in self.databases:
            return False
        
        # Remove from memory
        del self.databases[rag_name]
        self._md_snapshot_cache.pop(rag_name, None)
        
        # Remove from disk
        db_path = os.path.join(self.base_directory, rag_name)
        if os.path.exists(db_path):
            import shutil
            shutil.rmtree(db_path)
        
        return True

    @staticmethod
    def extract_invoice_fields(text: str) -> dict:
        """Extract structured invoice fields from OCR or PDF text using regex heuristics.

        Returns a dict with keys: invoice_no, date, vendor, total, gstin, pan.
        Each value is the matched string or None if not found.
        """
        def _first(pattern, flags=re.IGNORECASE):
            m = re.search(pattern, text, flags)
            return m.group(1).strip() if m else None

        return {
            "invoice_no": _first(
                r'invoice\s*(?:no\.?|number|#)\s*[:\-]?\s*([A-Z0-9\-/]+)'
            ),
            "date": _first(
                r'(?:invoice\s*date|date|dated)\s*[:\-]?\s*(\d{1,2}[./-]\d{1,2}[./-]\d{2,4})'
            ),
            "vendor": _first(
                r'^([A-Z][A-Za-z0-9\s&.,\-]{2,60}(?:Ltd\.?|Inc\.?|Pvt\.?|Corp\.?|LLP|LLC)?)',
                re.MULTILINE,
            ),
            "total": _first(
                r'(?:grand\s*total|total\s*amount|amount\s*due|net\s*payable|total)[:\s]*'
                r'[₹$€£]?\s*([\d,]+\.?\d*)'
            ),
            "gstin": _first(r'\b(GSTIN\s*[:\-]?\s*[A-Z0-9]{15})\b'),
            "pan": _first(r'\b([A-Z]{5}[0-9]{4}[A-Z])\b'),
        }
